import os
import sqlite3
from flask import Flask, render_template, request, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
import google.generativeai as genai
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
import fitz
import json
from datetime import datetime

# === Load Environment Variables ===
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
UPLOAD_FOLDER = "uploads"
DB = "buddy.db"

# === Configure Gemini ===
genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel("gemini-2.0-flash")

# === Flask App Setup ===
app = Flask(__name__)
CORS(app)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# === Database Initialization ===
def init_db():
    if os.path.exists(DB):
        try:
            with sqlite3.connect(DB) as con:
                con.execute("SELECT name FROM sqlite_master WHERE type='table';")
        except sqlite3.DatabaseError:
            print("⚠️ Corrupted DB detected. Reinitializing...")
            os.remove(DB)

    with sqlite3.connect(DB) as con:
        cur = con.cursor()
        cur.execute("""
            CREATE TABLE IF NOT EXISTS uploads (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                filename TEXT NOT NULL,
                filetype TEXT NOT NULL,
                filepath TEXT NOT NULL,
                extracted_text TEXT,
                uploaded_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS quiz (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                question TEXT NOT NULL,
                options TEXT NOT NULL,
                answer TEXT NOT NULL
            )
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS progress (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                related_upload_id INTEGER,
                action TEXT NOT NULL,
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        con.commit()

init_db()

# === Helper Function: Extract Text from Files ===
def extract_text(filepath, filetype):
    try:
        if filetype in ["png", "jpg", "jpeg"]:
            return pytesseract.image_to_string(Image.open(filepath))
        elif filetype == "txt":
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        elif filetype == "docx":
            return "\n".join([p.text for p in Document(filepath).paragraphs])
        elif filetype == "pptx":
            slides = []
            for slide in Presentation(filepath).slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides.append(shape.text)
            return "\n".join(slides)
        elif filetype == "pdf":
            text = ""
            with fitz.open(filepath) as doc:
                for page in doc:
                    text += page.get_text()
            return text
        return "Unsupported file type."
    except Exception as e:
        print(f"[Error] Text extraction failed: {e}")
        return f"Error extracting text: {str(e)}"

# === Routes ===
@app.route("/")
def index():
    return render_template("index.html")

@app.route("/upload", methods=["GET", "POST"])
def upload():
    if request.method == "POST":
        file = request.files.get("file")
        if not file or file.filename == "":
            return render_template("upload.html", message="No file selected")

        filename = f"{datetime.now().strftime('%Y%m%d%H%M%S')}_{secure_filename(file.filename)}"
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(filepath)

        filetype = filename.rsplit(".", 1)[1].lower()
        extracted = extract_text(filepath, filetype)

        with sqlite3.connect(DB) as con:
            cur = con.cursor()
            cur.execute("INSERT INTO uploads (filename, filetype, filepath, extracted_text) VALUES (?, ?, ?, ?)",
                        (filename, filetype, filepath, extracted))
            upload_id = cur.lastrowid
            cur.execute("INSERT INTO progress (related_upload_id, action) VALUES (?, ?)",
                        (upload_id, f"Uploaded {filename}"))
            con.commit()

        return render_template("upload.html", message="File uploaded successfully!")

    return render_template("upload.html")

@app.route("/library")
def library():
    with sqlite3.connect(DB) as con:
        cur = con.cursor()
        cur.execute("SELECT * FROM uploads ORDER BY uploaded_at DESC")
        files = cur.fetchall()
    return render_template("library.html", files=files)

@app.route("/progress")
def progress():
    with sqlite3.connect(DB) as con:
        cur = con.cursor()
        cur.execute("""
            SELECT p.id, p.action, p.timestamp, u.filename
            FROM progress p
            LEFT JOIN uploads u ON p.related_upload_id = u.id
            ORDER BY p.timestamp DESC
        """)
        logs = cur.fetchall()
    return render_template("progress.html", logs=logs)

@app.route("/quiz", methods=["GET"])
def quiz():
    questions = []
    with sqlite3.connect(DB) as con:
        cur = con.cursor()
        cur.execute("SELECT id, question, options, answer FROM quiz")
        rows = cur.fetchall()
        for row in rows:
            options_list = row[2].split(";") if row[2] else []
            questions.append({
                "id": row[0],
                "question_text": row[1],
                "options": options_list,
                "correct_answer": row[3]
            })
    return render_template("quiz.html", questions=questions)

@app.route("/generate-quiz", methods=["POST"])
def generate_quiz():
    try:
        with sqlite3.connect(DB) as con:
            cur = con.cursor()
            cur.execute("SELECT extracted_text FROM uploads")
            texts = [row[0] for row in cur.fetchall() if row[0]]

        context = "\n".join(texts)
        if not context:
            return jsonify({"success": False, "message": "No study material uploaded."})

        quiz_prompt = f'''
        Based on the following study material, generate 3 multiple-choice questions.
        Each question must include a question string, exactly 4 options, and the correct answer.

        Return only valid JSON array format like:
        [
          {{"question": "What is ...?", "options": ["A", "B", "C", "D"], "answer": "A"}}
        ]
        Study Material:
        {context}
        '''

        response = model.generate_content(quiz_prompt, generation_config={"temperature": 0.7, "response_mime_type": "application/json"})
        quiz_data = response.text

        try:
            generated_questions_list = json.loads(quiz_data)
        except json.JSONDecodeError:
            return jsonify({"success": False, "message": "AI response could not be parsed."})

        with sqlite3.connect(DB) as con:
            cur = con.cursor()
            cur.execute("DELETE FROM quiz")
            for q in generated_questions_list:
                question_text = q.get("question", "").strip()
                options_text = ";".join(q.get("options", []))
                answer_text = q.get("answer", "").strip()

                if question_text and options_text and answer_text:
                    cur.execute(
                        "INSERT INTO quiz (question, options, answer) VALUES (?, ?, ?)",
                        (question_text, options_text, answer_text)
                    )
            cur.execute("INSERT INTO progress (action) VALUES (?)", ("Generated quiz",))
            con.commit()

        return jsonify({"success": True, "message": "Quiz generated."})

    except Exception as e:
        print("[Quiz Error]", str(e))
        return jsonify({"success": False, "message": str(e)})

@app.route("/ask-buddy", methods=["POST"])
def ask_buddy():
    question = request.get_json().get("question", "")
    with sqlite3.connect(DB) as con:
        cur = con.cursor()
        cur.execute("SELECT extracted_text FROM uploads")
        context = "\n".join([row[0] for row in cur.fetchall() if row[0]])

    prompt = f"Use this material to answer: {question}\n\nContext:\n{context}"
    try:
        response = model.generate_content(prompt)
        answer = response.text
    except Exception as e:
        answer = f"Error: {str(e)}"

    with sqlite3.connect(DB) as con:
        cur = con.cursor()
        cur.execute("INSERT INTO progress (action) VALUES (?)", (f"Asked: {question}",))
        con.commit()

    return jsonify({"answer": answer})

@app.route("/about")
def about():
    return render_template("about.html")

@app.route("/analyze-complexity", methods=["POST"])
def analyze_complexity():
    code = request.json.get("code", "")
    prompt = f"""
    Analyze the following code and provide:
    - Time complexity (worst-case)
    - Space complexity
    - A short explanation (3–4 lines)
    Return JSON:
    {{ "time": "...", "space": "...", "explanation": "..." }}
    Code:
    {code}
    """
    try:
        response = model.generate_content(prompt)
        return jsonify(json.loads(response.text))
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True)
